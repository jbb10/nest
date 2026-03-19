#!/usr/bin/env python3
"""Generate minimal test fixtures for E2E tests.

Creates small test documents (<100KB each) for PDF, DOCX, PPTX, and XLSX.
Also creates a corrupt PDF for negative testing, and a PDF with embedded
raster images for vision pipeline E2E tests (Story 7.5).

Run: uv run python tests/e2e/fixtures/generate_fixtures.py
"""

from pathlib import Path

# DOCX generation
from docx import Document
from docx.shared import Pt

# XLSX generation
from openpyxl import Workbook

# PPTX generation
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt

# PDF generation using reportlab
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


def create_pdf(output_path: Path) -> None:
    """Create a minimal PDF with title and paragraphs."""
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter

    # Title
    c.setFont("Helvetica-Bold", 24)
    c.drawString(72, height - 72, "Quarterly Report Q4 2025")

    # Paragraphs
    c.setFont("Helvetica", 12)
    y = height - 120

    paragraphs = [
        "Executive Summary",
        "This quarterly report provides an overview of our performance during Q4 2025.",
        "Key highlights include revenue growth of 15% and expansion into new markets.",
        "",
        "Financial Highlights",
        "Revenue: $12.5M (up 15% YoY)",
        "Operating Margin: 22%",
        "Customer Acquisition: 500 new enterprise clients",
        "",
        "Strategic Initiatives",
        "We have successfully launched three new product lines and established",
        "partnerships with leading technology providers in the EMEA region.",
    ]

    for para in paragraphs:
        c.drawString(72, y, para)
        y -= 18

    c.save()
    print(f"Created: {output_path} ({output_path.stat().st_size} bytes)")


def create_docx(output_path: Path) -> None:
    """Create a minimal DOCX with title, bullets, and paragraph."""
    doc = Document()

    # Title
    title = doc.add_heading("Summary Document", level=1)
    title.style.font.size = Pt(24)

    # Introduction paragraph
    doc.add_paragraph(
        "This document summarizes the key findings from our recent analysis. "
        "The data shows promising trends across all major metrics."
    )

    # Bullet points
    doc.add_heading("Key Points", level=2)
    bullets = [
        "Market share increased by 8% in Q4",
        "Customer satisfaction score improved to 94%",
        "Product delivery time reduced by 20%",
        "New feature adoption rate exceeded targets",
    ]
    for bullet in bullets:
        doc.add_paragraph(bullet, style="List Bullet")

    # Conclusion
    doc.add_heading("Conclusion", level=2)
    doc.add_paragraph(
        "Based on these findings, we recommend continuing the current strategy "
        "while exploring additional opportunities in emerging markets."
    )

    doc.save(output_path)
    print(f"Created: {output_path} ({output_path.stat().st_size} bytes)")


def create_pptx(output_path: Path) -> None:
    """Create a minimal PPTX with 3 slides."""
    prs = Presentation()

    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Q4 2025 Presentation"
    tf.paragraphs[0].font.size = PptPt(44)
    tf.paragraphs[0].font.bold = True

    # Slide 2: Bullet points
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Key Metrics"
    tf.paragraphs[0].font.size = PptPt(32)
    tf.paragraphs[0].font.bold = True

    bullets = ["Revenue: $12.5M", "Growth: 15%", "Customers: 2,500+"]
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = PptPt(24)
        p.level = 0

    # Slide 3: Summary
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.text = "Thank You"
    tf.paragraphs[0].font.size = PptPt(44)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = "Questions? Contact team@example.com"
    p.font.size = PptPt(20)

    prs.save(output_path)
    print(f"Created: {output_path} ({output_path.stat().st_size} bytes)")


def create_xlsx(output_path: Path) -> None:
    """Create a minimal XLSX with headers and data rows."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales Data"

    # Headers
    headers = ["Region", "Q1 Sales", "Q2 Sales", "Q3 Sales", "Q4 Sales", "Total"]
    ws.append(headers)

    # Data rows
    data = [
        ["North America", 2500, 2800, 3100, 3500, 11900],
        ["Europe", 1800, 2100, 2300, 2700, 8900],
        ["Asia Pacific", 1200, 1500, 1800, 2200, 6700],
        ["Latin America", 800, 900, 1100, 1400, 4200],
        ["Middle East", 400, 500, 600, 800, 2300],
    ]
    for row in data:
        ws.append(row)

    # Totals row
    ws.append(["Total", 6700, 7800, 8900, 10600, 34000])

    wb.save(output_path)
    print(f"Created: {output_path} ({output_path.stat().st_size} bytes)")


def create_corrupt_pdf(output_path: Path, source_pdf: Path) -> None:
    """Create a corrupt PDF by truncating a valid one."""
    with open(source_pdf, "rb") as f:
        corrupt_bytes = f.read(100)  # Just first 100 bytes

    with open(output_path, "wb") as f:
        f.write(corrupt_bytes)

    print(f"Created: {output_path} ({output_path.stat().st_size} bytes)")


def create_image_pdf(output_path: Path) -> None:
    """Create a PDF with embedded raster images for vision pipeline E2E tests.

    Uses PIL to generate a bar-chart image and embeds it in a PDF via
    reportlab. PIL is available as a transitive dependency of docling.

    The resulting PDF should be detected by Docling's layout analysis
    as containing at least one figure/picture element.
    """
    from io import BytesIO

    from PIL import Image, ImageDraw
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas as rl_canvas

    def _make_bar_chart_image() -> BytesIO:
        """Create a colourful bar chart image (400×280 px) using PIL."""
        img = Image.new("RGB", (400, 280), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Outer border
        draw.rectangle([0, 0, 399, 279], outline=(50, 50, 50), width=2)

        # Title banner
        draw.rectangle([0, 0, 399, 36], fill=(70, 130, 180))
        draw.text((120, 10), "Quarterly Sales 2025", fill=(255, 255, 255))

        # Grid lines and Y-axis labels
        base_y = 245
        for val, y_off in [(200, 200), (150, 150), (100, 100), (50, 50)]:
            y_pos = base_y - y_off
            draw.line([45, y_pos, 390, y_pos], fill=(210, 210, 210), width=1)
            draw.text((5, y_pos - 7), str(val), fill=(90, 90, 90))

        # Bars: (label, bar_height, colour)
        bars = [
            ("Q1", 120, (70, 130, 180)),
            ("Q2", 165, (255, 165, 0)),
            ("Q3", 205, (60, 179, 113)),
            ("Q4", 180, (220, 20, 60)),
        ]
        bar_width = 55
        for i, (label, bh, colour) in enumerate(bars):
            x = 55 + i * 82
            draw.rectangle([x, base_y - bh, x + bar_width, base_y], fill=colour)
            draw.text((x + 15, base_y + 5), label, fill=(50, 50, 50))

        # X-axis baseline
        draw.line([45, base_y, 390, base_y], fill=(50, 50, 50), width=2)

        buf = BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf

    c = rl_canvas.Canvas(str(output_path), pagesize=letter)
    page_width, page_height = letter

    # Page heading
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, page_height - 72, "Technical Analysis Report with Figures")

    # Figure caption
    c.setFont("Helvetica", 11)
    c.drawString(72, page_height - 104, "Figure 1: Quarterly Sales Performance (2025)")

    # Embed the bar-chart image
    chart_buf = _make_bar_chart_image()
    ir = ImageReader(chart_buf)
    img_w, img_h = 380, 266  # slightly scaled from 400×280
    c.drawImage(ir, 72, page_height - 104 - img_h - 8, width=img_w, height=img_h)

    # Body text below the figure
    c.setFont("Helvetica", 10)
    body_y = page_height - 104 - img_h - 32
    for line in [
        "The bar chart above summarises quarterly revenue for 2025.",
        "Q3 achieved peak performance with sustained growth in all regions.",
        "Data sourced from the consolidated financial reporting system.",
    ]:
        c.drawString(72, body_y, line)
        body_y -= 16

    c.save()
    print(f"Created: {output_path} ({output_path.stat().st_size} bytes)")


def main() -> None:
    """Generate all test fixtures."""
    fixtures_dir = Path(__file__).parent

    # Create main test documents
    pdf_path = fixtures_dir / "quarterly.pdf"
    create_pdf(pdf_path)
    create_docx(fixtures_dir / "summary.docx")
    create_pptx(fixtures_dir / "deck.pptx")
    create_xlsx(fixtures_dir / "data.xlsx")

    # Create corrupt PDF for negative tests
    create_corrupt_pdf(fixtures_dir / "corrupt.pdf", pdf_path)

    # Create PDF with embedded images for vision pipeline tests (Story 7.5)
    create_image_pdf(fixtures_dir / "image_doc.pdf")

    print("\n✅ All fixtures generated successfully!")


if __name__ == "__main__":
    main()
